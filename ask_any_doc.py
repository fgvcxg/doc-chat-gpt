import os
from langchain.schema import Document
from langchain.vectorstores import Chroma
from langchain.embeddings import FakeEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from langchain.document_loaders import PyPDFLoader
from pptx import Presentation
import pandas as pd
import docx

# 대상 폴더
doc_folder = "D:/toyproj"
all_docs = []

print("📂 문서 로딩 시작...")

# ✅ 1. PDF
for file in os.listdir(doc_folder):
    if file.endswith(".pdf"):
        loader = PyPDFLoader(os.path.join(doc_folder, file))
        all_docs.extend(loader.load())

# ✅ 2. PPTX
for file in os.listdir(doc_folder):
    if file.endswith(".pptx"):
        path = os.path.join(doc_folder, file)
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        all_docs.append(Document(page_content=text, metadata={"source": file}))

# ✅ 3. XLSX
for file in os.listdir(doc_folder):
    if file.endswith(".xlsx"):
        path = os.path.join(doc_folder, file)
        xls = pd.read_excel(path, sheet_name=None)  # 모든 시트 불러오기
        text = ""
        for sheet_name, df in xls.items():
            text += f"▶ 시트: {sheet_name}\n"
            text += df.to_string(index=False) + "\n\n"
        all_docs.append(Document(page_content=text, metadata={"source": file}))

# ✅ 4. DOCX
for file in os.listdir(doc_folder):
    if file.endswith(".docx"):
        path = os.path.join(doc_folder, file)
        doc = docx.Document(path)
        text = "\n".join([p.text for p in doc.paragraphs])
        all_docs.append(Document(page_content=text, metadata={"source": file}))

print(f"✅ 총 문서 수: {len(all_docs)}개")

# 텍스트 분할
splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
split_docs = splitter.split_documents(all_docs)

# GPT 없이 FAKE 임베딩으로 검색엔진 구성
embedding = FakeEmbeddings(size=768)
db = Chroma.from_documents(split_docs, embedding)
retriever = db.as_retriever()

# 질의 응답 루프
print("\n🔍 문서 검색 시스템 시작 (종료하려면 'exit')")
while True:
    query = input("\n질문 키워드: ")
    if query.lower() in ["exit", "quit"]:
        print("👋 종료합니다. 감사합니다!")
        break

    results = retriever.get_relevant_documents(query)

    if not results:
        print("❌ 관련 내용을 찾을 수 없습니다.")
    else:
        print(f"\n📌 관련 문서 조각 {len(results)}개:")
        for i, doc in enumerate(results):
            print(f"\n--- 문서 {i+1} ({doc.metadata['source']}) ---")
            print(doc.page_content.strip()[:500] + "...")
