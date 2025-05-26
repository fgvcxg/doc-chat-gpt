import os
import streamlit as st
from dotenv import load_dotenv
from langchain.schema import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain.chains import RetrievalQA
from pptx import Presentation
import pandas as pd
import docx

# ✅ .env에서 API 키 불러오기
load_dotenv()
api_key = os.getenv("OPENAI_API_KEY")

# ✅ Streamlit 기본 설정
st.set_page_config(page_title="문서 기반 GPT 챗봇", layout="wide")
st.title("📚 문서 기반 대화형 GPT 챗봇")

# 앱 실행 디렉토리 기준 상대 경로
doc_dir = "docs"

# 없으면 자동 생성 (Streamlit Cloud 대응)
if not os.path.exists(doc_dir):
    os.makedirs(doc_dir)

# ✅ 문서 불러오기 버튼
if st.button("✅ 문서 로드 및 챗봇 준비"):
    all_docs = []
    st.info("📁 문서를 불러오는 중...")

    for file in os.listdir(doc_dir):
        path = os.path.join(doc_dir, file)
        if file.endswith(".pdf"):
            loader = PyPDFLoader(path)
            all_docs.extend(loader.load())
        elif file.endswith(".pptx"):
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            all_docs.append(Document(page_content=text, metadata={"source": file}))
        elif file.endswith(".xlsx"):
            xls = pd.read_excel(path, sheet_name=None)
            text = ""
            for sheet_name, df in xls.items():
                text += f"[{sheet_name}]\n" + df.to_string(index=False) + "\n\n"
            all_docs.append(Document(page_content=text, metadata={"source": file}))
        elif file.endswith(".docx"):
            doc = docx.Document(path)
            text = "\n".join([p.text for p in doc.paragraphs])
            all_docs.append(Document(page_content=text, metadata={"source": file}))

    # ✅ 문서 조각 및 벡터화
    splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    split_docs = splitter.split_documents(all_docs)
    vectordb = FAISS.from_documents(split_docs, OpenAIEmbeddings(openai_api_key=api_key))
    retriever = vectordb.as_retriever()

    # ✅ GPT-3.5 연결
    qa_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(model="gpt-3.5-turbo", openai_api_key=api_key),
        retriever=retriever,
        return_source_documents=True
    )
    st.session_state.qa_chain = qa_chain
    st.success("문서 로딩 및 챗봇 준비 완료!")

# ✅ 질문 입력창 (대화형 스타일)
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

user_input = st.chat_input("질문을 입력하세요")

if user_input and "qa_chain" in st.session_state:
    with st.spinner("GPT가 답변 중입니다..."):
        result = st.session_state.qa_chain({"query": user_input})
        answer = result["result"]
        sources = result["source_documents"]

        # 💬 대화 내용 저장
        st.session_state.chat_history.append(("user", user_input))
        st.session_state.chat_history.append(("gpt", answer))

        # 🧠 대화 내용 출력
        for role, msg in st.session_state.chat_history:
            if role == "user":
                st.chat_message("user").write(msg)  # 오른쪽 정렬
            else:
                st.chat_message("assistant").write(msg)  # 왼쪽 정렬

        # 📁 참조 문서 표시
        with st.expander("📎 참조 문서 보기"):
            for doc in sources:
                st.markdown(f"📄 **{doc.metadata['source']}**")
                st.code(doc.page_content[:500] + "...")
